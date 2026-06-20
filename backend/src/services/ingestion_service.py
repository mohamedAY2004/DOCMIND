"""Parse uploaded files into text chunks.

Supports PDF (via PyMuPDF), PPTX (via python-pptx), PNG (OCR not yet
implemented — returns a descriptive stub), and plain text / markdown files.

Chunking strategy
-----------------
The goal is good retrieval over college material (textbooks, slide decks,
lecture PDFs). Three things matter here and drive the design:

1. **Tables stay intact.** Tables are detected with PyMuPDF's native
   ``page.find_tables()`` (PDF) / ``shape.has_table`` (PPTX), rendered to
   Markdown, and emitted as *atomic* chunks (``metadata["kind"] == "table"``).
   They are never sliced and their text is excluded from the prose stream so it
   is not duplicated as garbled run-on text.

2. **Structure-aware splitting.** Prose is split on paragraph → line →
   sentence → word boundaries (``_split``) instead of hard char slicing, so a
   definition or worked example is far less likely to be cut mid-thought.

3. **Heading breadcrumbs.** Headings are detected by font size and prepended to
   each chunk (``[Section Title]\n…``) and stored in ``metadata["section"]``,
   so a chunk that only says "this eliminates transitive dependencies" still
   carries the "Third Normal Form" context it needs to be retrievable.

The public surface (``IngestedChunk``, ``ingest_file``, ``ingest_pdf``,
``ingest_pptx``, ``ingest_png``, ``detect_pdf_encrypted``) is unchanged;
downstream services only read ``chunk.text`` / ``chunk.metadata`` and treat the
metadata dict opaquely, so the extra ``kind``/``section`` keys are additive.
"""
from __future__ import annotations

import logging
import re
from dataclasses import dataclass
from pathlib import Path
from typing import List, Optional, Tuple

logger = logging.getLogger("docmind.ingestion")

# Char-based sizing keeps us dependency-free (no tokenizer). ~1000 chars is
# roughly 250 tokens, comfortably within any embedding model's context while
# keeping related material together.
DEFAULT_CHUNK_SIZE = 1000
DEFAULT_OVERLAP = 150


@dataclass
class IngestedChunk:
    text: str
    metadata: dict


# --------------------------------------------------------------------------- #
# Text cleaning + structure-aware splitting
# --------------------------------------------------------------------------- #
def _clean(text: str) -> str:
    """Light cleanup that preserves paragraph structure.

    Unlike a naive cleaner we do **not** flatten single newlines into spaces —
    paragraph/line structure is what the splitter keys off of. We only drop
    bare page-number lines, de-hyphenate words broken across lines, and
    collapse runs of spaces / blank lines.
    """
    text = re.sub(r"^\s*\d+\s*$", "", text, flags=re.MULTILINE)  # bare page nums
    text = re.sub(r"-\n(\w)", r"\1", text)                       # de-hyphenate
    text = re.sub(r"[ \t]{2,}", " ", text)
    text = re.sub(r"\n{3,}", "\n\n", text)
    return text.strip()


_PARA_RE = re.compile(r"\n\s*\n")
_SENT_RE = re.compile(r"(?<=[.!?])\s+")


def _atomize(text: str, max_size: int) -> List[str]:
    """Break text into the largest readable units that fit ``max_size``.

    Recursively backs off paragraph → line → sentence → hard-slice, so atoms
    respect natural boundaries wherever possible (LangChain-style recursion,
    no dependency).
    """
    atoms: List[str] = []
    for para in _PARA_RE.split(text):
        para = para.strip()
        if not para:
            continue
        if len(para) <= max_size:
            atoms.append(para)
            continue
        for line in para.split("\n"):
            line = line.strip()
            if not line:
                continue
            if len(line) <= max_size:
                atoms.append(line)
                continue
            for sent in _SENT_RE.split(line):
                sent = sent.strip()
                if not sent:
                    continue
                if len(sent) <= max_size:
                    atoms.append(sent)
                else:  # pathological single token — hard slice as last resort
                    for i in range(0, len(sent), max_size):
                        atoms.append(sent[i:i + max_size])
    return atoms


def _merge_with_overlap(atoms: List[str], chunk_size: int, overlap: int) -> List[str]:
    """Greedily pack atoms up to ``chunk_size`` with a trailing-atom overlap."""
    chunks: List[str] = []
    cur: List[str] = []
    cur_len = 0
    for atom in atoms:
        add = len(atom) + (1 if cur else 0)
        if cur and cur_len + add > chunk_size:
            chunks.append(" ".join(cur))
            # Carry trailing atoms (~overlap chars) into the next chunk.
            tail: List[str] = []
            tail_len = 0
            for a in reversed(cur):
                if tail_len + len(a) + 1 > overlap:
                    break
                tail.insert(0, a)
                tail_len += len(a) + 1
            cur = tail
            cur_len = sum(len(a) + 1 for a in cur)
        cur.append(atom)
        cur_len += add
    if cur:
        chunks.append(" ".join(cur))
    return chunks


def _split(text: str, chunk_size: int = DEFAULT_CHUNK_SIZE,
           overlap: int = DEFAULT_OVERLAP) -> List[str]:
    """Structure-aware splitter with overlap. No external dependency."""
    if not text or not text.strip():
        return []
    return _merge_with_overlap(_atomize(text, chunk_size), chunk_size, overlap)


def _with_breadcrumb(text: str, heading: Optional[str]) -> str:
    """Prepend the section heading so the chunk is self-describing."""
    if heading and not text.lstrip().startswith(f"[{heading}]"):
        return f"[{heading}]\n{text}"
    return text


# --------------------------------------------------------------------------- #
# PDF: native table extraction + font-size heading detection
# --------------------------------------------------------------------------- #
def _mode(values: List[float]) -> float:
    return Counter(values).most_common(1)[0][0] if values else 0.0


def _mode(values: List[float]) -> float:
    return max(set(values), key=values.count) if values else 0.0


def _intersects(a, b) -> bool:
    """Axis-aligned bbox intersection for (x0, y0, x1, y1) tuples."""
    return not (a[2] <= b[0] or b[2] <= a[0] or a[3] <= b[1] or b[3] <= a[1])


def _intersects_any(bbox, boxes) -> bool:
    return bbox is not None and any(_intersects(bbox, b) for b in boxes)


def _rows_to_md(rows: List[list]) -> str:
    """Fallback Markdown table renderer when ``to_markdown`` is unavailable."""
    def cell(c) -> str:
        return ("" if c is None else str(c)).strip().replace("\n", " ").replace("|", "\\|")

    if not rows:
        return ""
    ncols = max(len(r) for r in rows)
    out = ["| " + " | ".join(cell(c) for c in (rows[0] + [""] * (ncols - len(rows[0])))) + " |"]
    out.append("| " + " | ".join(["---"] * ncols) + " |")
    for r in rows[1:]:
        out.append("| " + " | ".join(cell(c) for c in (r + [""] * (ncols - len(r)))) + " |")
    return "\n".join(out)


def _extract_tables(page) -> Tuple[list, List[str]]:
    """Return (table bboxes, table Markdown strings) for a PDF page."""
    boxes: list = []
    mds: List[str] = []
    try:
        finder = page.find_tables()
    except Exception:  # noqa: BLE001 — find_tables can raise on odd pages
        return boxes, mds
    for t in getattr(finder, "tables", []):
        try:
            data = t.extract()
        except Exception:  # noqa: BLE001
            continue
        nrows = len(data)
        ncols = max((len(r) for r in data), default=0)
        if nrows < 2 or ncols < 2:  # skip false positives (single col/row)
            continue
        try:
            md = (t.to_markdown() or "").strip()
        except Exception:  # noqa: BLE001
            md = _rows_to_md(data)
        if md:
            mds.append(md)
            boxes.append(tuple(t.bbox))
    return boxes, mds


def _join_lines(lines: List[str]) -> str:
    """Join a block's lines, de-hyphenating words split across lines."""
    out = ""
    for ln in lines:
        if not out:
            out = ln
        elif out.endswith("-") and len(out) >= 2 and out[-2].isalpha():
            out = out[:-1] + ln
        else:
            out += " " + ln
    return out.strip()


def _prose_blocks(page, table_boxes) -> Tuple[List[List[Tuple[str, float]]], float]:
    """Extract non-table text as blocks of (line_text, line_font_size).

    Heading detection works at line granularity (below) rather than block
    granularity, because PDF producers often pack a heading and the paragraph
    that follows it into a single text block.
    """
    try:
        data = page.get_text("dict")
    except Exception:  # noqa: BLE001
        return [], 0.0
    # Body font = the size carrying the most *characters*, not the most lines.
    # Weighting by character mass is robust when a page has few body lines.
    size_chars: dict = {}
    blocks: List[List[Tuple[str, float]]] = []
    for block in data.get("blocks", []):
        if block.get("type") != 0:  # 0 == text block; skip images
            continue
        if _intersects_any(block.get("bbox"), table_boxes):
            continue
        lines: List[Tuple[str, float]] = []
        for line in block.get("lines", []):
            spans = line.get("spans", [])
            txt = "".join(s.get("text", "") for s in spans).strip()
            if not txt:
                continue
            line_size = 0.0
            for s in spans:
                sz = round(float(s.get("size", 0.0)), 1)
                line_size = max(line_size, sz)
                stripped = s.get("text", "").strip()
                if stripped:
                    size_chars[sz] = size_chars.get(sz, 0) + len(stripped)
            lines.append((txt, line_size))
        if lines:
            blocks.append(lines)
    body_size = max(size_chars, key=size_chars.get) if size_chars else 0.0
    return blocks, body_size


def _looks_like_heading(text: str, size: float, body_size: float) -> bool:
    return bool(
        body_size
        and text
        and len(text) <= 120
        and size >= body_size * 1.18
        and not text.endswith((".", "!", "?", ",", ";"))
    )


def _segments_from_blocks(blocks, body_size, heading):
    """Turn blocks into (paragraph, heading) pairs, tracking the current heading.

    ``heading`` is threaded in and out so a section that continues across a page
    break keeps its breadcrumb.
    """
    segments: List[Tuple[str, Optional[str]]] = []
    for lines in blocks:
        para: List[str] = []
        for txt, size in lines:
            if _looks_like_heading(txt, size, body_size):
                if para:
                    segments.append((_join_lines(para), heading))
                    para = []
                heading = txt
            else:
                para.append(txt)
        if para:
            segments.append((_join_lines(para), heading))
    return segments, heading


def _group_segments(segments):
    """Merge consecutive paragraphs that share a heading into one block."""
    groups: List[dict] = []
    for text, heading in segments:
        if groups and groups[-1]["heading"] == heading:
            groups[-1]["paras"].append(text)
        else:
            groups.append({"heading": heading, "paras": [text]})
    return [{"heading": g["heading"], "text": "\n\n".join(g["paras"])} for g in groups]


def ingest_pdf(path: Path) -> List[IngestedChunk]:
    try:
        import fitz  # PyMuPDF
    except ImportError:
        logger.error("PyMuPDF not installed; cannot ingest PDF.")
        return []

    chunks: List[IngestedChunk] = []
    heading: Optional[str] = None  # carried across pages
    with fitz.open(str(path)) as doc:
        for page_idx, page in enumerate(doc):
            page_no = page_idx + 1
            part = 0

            table_boxes, table_mds = _extract_tables(page)
            for md in table_mds:
                chunks.append(
                    IngestedChunk(
                        text=_with_breadcrumb(md, heading),
                        metadata={
                            "source": path.name,
                            "page": page_no,
                            "part": part,
                            "kind": "table",
                            "section": heading,
                        },
                    )
                )
                part += 1

            blocks_info, body_size = _prose_blocks(page, table_boxes)
            segments, heading = _segments_from_blocks(blocks_info, body_size, heading)
            for grp in _group_segments(segments):
                cleaned = _clean(grp["text"])
                for piece in _split(cleaned):
                    chunks.append(
                        IngestedChunk(
                            text=_with_breadcrumb(piece, grp["heading"]),
                            metadata={
                                "source": path.name,
                                "page": page_no,
                                "part": part,
                                "kind": "text",
                                "section": grp["heading"],
                            },
                        )
                    )
                    part += 1
    return chunks


# --------------------------------------------------------------------------- #
# PPTX: table-aware, slide title as heading
# --------------------------------------------------------------------------- #
def _pptx_table_to_md(table) -> str:
    rows: List[list] = []
    for row in table.rows:
        rows.append([cell.text for cell in row.cells])
    return _rows_to_md(rows)


def _slide_title(slide) -> Optional[str]:
    try:
        title = slide.shapes.title
        if title is not None and title.text and title.text.strip():
            return title.text.strip()
    except Exception:  # noqa: BLE001
        pass
    return None


def ingest_pptx(path: Path) -> List[IngestedChunk]:
    try:
        from pptx import Presentation
    except ImportError:
        logger.error("python-pptx not installed; cannot ingest PPTX.")
        return []

    chunks: List[IngestedChunk] = []
    prs = Presentation(str(path))
    for slide_idx, slide in enumerate(prs.slides):
        slide_no = slide_idx + 1
        part = 0
        title = _slide_title(slide)
        try:
            title_shape = slide.shapes.title
        except Exception:  # noqa: BLE001
            title_shape = None

        buffer: List[str] = []
        for shape in slide.shapes:
            if shape is title_shape:
                continue  # title becomes the breadcrumb, not body text
            if getattr(shape, "has_table", False):
                md = _pptx_table_to_md(shape.table)
                if md:
                    chunks.append(
                        IngestedChunk(
                            text=_with_breadcrumb(md, title),
                            metadata={
                                "source": path.name,
                                "slide": slide_no,
                                "part": part,
                                "kind": "table",
                                "section": title,
                            },
                        )
                    )
                    part += 1
            elif getattr(shape, "has_text_frame", False) and shape.text and shape.text.strip():
                buffer.append(shape.text.strip())

        cleaned = _clean("\n\n".join(buffer))
        for piece in _split(cleaned):
            chunks.append(
                IngestedChunk(
                    text=_with_breadcrumb(piece, title),
                    metadata={
                        "source": path.name,
                        "slide": slide_no,
                        "part": part,
                        "kind": "text",
                        "section": title,
                    },
                )
            )
            part += 1
    return chunks


def ingest_png(path: Path) -> List[IngestedChunk]:
    """OCR hook. Deferred — returns a stub so the pipeline still completes."""
    logger.warning("PNG OCR is not implemented; returning placeholder chunk for %s", path)
    return [
        IngestedChunk(
            text=f"[Image upload: {path.name}. Visual content has not yet been transcribed.]",
            metadata={"source": path.name, "kind": "png_placeholder"},
        )
    ]


def ingest_file(path: Path) -> List[IngestedChunk]:
    """Dispatch based on extension."""
    ext = path.suffix.lower()
    if ext == ".pdf":
        return ingest_pdf(path)
    if ext == ".pptx":
        return ingest_pptx(path)
    if ext == ".png":
        return ingest_png(path)
    if ext in {".txt", ".md"}:
        text = path.read_text(encoding="utf-8", errors="replace")
        return [
            IngestedChunk(
                text=piece,
                metadata={"source": path.name, "part": i, "kind": "text"},
            )
            for i, piece in enumerate(_split(_clean(text)))
        ]
    logger.warning("Unsupported extension %s for %s", ext, path)
    return []


def detect_pdf_encrypted(path: Path) -> bool:
    """Return True only if the PDF is actually locked behind a user password.

    Many PDFs produced by Word, Google Docs, macOS Preview, and "print to PDF"
    pipelines carry an owner-permissions dictionary that ``pypdf`` reports as
    ``is_encrypted = True`` even though no password is required to read them.
    We attempt to authenticate with an empty password first and only treat the
    file as encrypted when that attempt fails.
    """
    try:
        import fitz  # PyMuPDF
    except ImportError:
        fitz = None

    if fitz is not None:
        try:
            with fitz.open(str(path)) as doc:
                if not doc.needs_pass:
                    return False
                return not doc.authenticate("")
        except Exception:
            return False

    try:
        from pypdf import PdfReader
    except ImportError:
        return False
    try:
        reader = PdfReader(str(path))
        if not reader.is_encrypted:
            return False
        try:
            result = reader.decrypt("")
        except Exception:
            return True
        return not bool(result)
    except Exception:
        return False
